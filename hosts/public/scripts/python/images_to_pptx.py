from pptx import Presentation
from pptx.util import Inches
import os
import argparse
import subprocess
import time
from pptx.enum.action import PP_ACTION
from pptx.dml.color import RGBColor
import shutil


def create_presentation_from_images(image_folder, output_file, add_thumbnail=True):
    """
    Create a PowerPoint presentation from a folder of images.
    Each image will be placed on a separate slide.
    """
    # Create a blank presentation
    prs = Presentation()

    # Define slide width and height for 1920x1080 (16:9) presentation
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # Get all image files from the folder
    image_files = [
        f for f in os.listdir(image_folder)
        if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp'))
    ]

    # Sort the image files to maintain order
    image_files.sort()

    if not image_files:
        print(f"No image files found in {image_folder}")
        return

    # For each image, create a slide
    for image_file in image_files:
        # Add a blank slide
        slide_layout = prs.slide_layouts[6]  # 6 is for blank slide
        slide = prs.slides.add_slide(slide_layout)

        # Add the image to fill the slide
        img_path = os.path.join(image_folder, image_file)

                # Calculate dimensions to maintain aspect ratio
        from PIL import Image
        
        with Image.open(img_path) as img:
            img_width, img_height = img.size
            
            # Calculate ratios
            width_ratio = prs.slide_width / img_width
            height_ratio = prs.slide_height / img_height
            
            # Use the smaller ratio to ensure the image fits within the slide
            ratio = min(width_ratio, height_ratio)
            
            new_width = img_width * ratio
            new_height = img_height * ratio
            
            # Calculate position to center the image on the slide
            left = (prs.slide_width - new_width) / 2
            top = (prs.slide_height - new_height) / 2
            
            # Add picture to slide with calculated dimensions
            slide.shapes.add_picture(img_path, left, top, width=new_width, height=new_height)
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)  # RGB for black



        print(f"Added {image_file} to presentation")

    # Add document properties
    core_properties = prs.core_properties
    core_properties.title = "Image Presentation"
    core_properties.subject = "Created with python-pptx"
    core_properties.author = "Automatic Generator"
    core_properties.comments = "Generated from image files"

    # Save the presentation
    temp_output = output_file
    prs.save(temp_output)
    print(f"Presentation saved as {temp_output}")
    
    if add_thumbnail:
        # Use LibreOffice to process the file to generate proper thumbnail
        try:
            process_with_libreoffice(temp_output)
        except Exception as e:
            print(f"Error generating thumbnail: {e}")


def process_with_libreoffice(pptx_file):
    """
    Process the PPTX file with LibreOffice to generate proper thumbnails.
    This opens and saves the file again using LibreOffice.
    """
    # Create a temp directory for conversion
    temp_dir = os.path.splitext(pptx_file)[0] + "_temp"
    os.makedirs(temp_dir, exist_ok=True)
    
    # Copy the file to temp dir
    temp_file = os.path.join(temp_dir, os.path.basename(pptx_file))
    shutil.copy2(pptx_file, temp_file)
    
    try:
        # Use soffice (LibreOffice) to open and export the file
        command = [
            "soffice",
            "--headless",
            "--convert-to", "pptx",
            "--outdir", os.path.dirname(pptx_file),
            temp_file
        ]
        
        print("Running LibreOffice to process the presentation...")
        process = subprocess.Popen(command, 
                                  stdout=subprocess.PIPE, 
                                  stderr=subprocess.PIPE)
        stdout, stderr = process.communicate(timeout=30)
        
        if process.returncode != 0:
            print(f"LibreOffice error: {stderr.decode()}")
        else:
            print("LibreOffice processing completed successfully")
            
            # If the command succeeded but didn't overwrite the original
            processed_file = os.path.join(
                os.path.dirname(pptx_file), 
                os.path.basename(temp_file)
            )
            if os.path.exists(processed_file) and processed_file != pptx_file:
                shutil.move(processed_file, pptx_file)
                print(f"Moved processed file to {pptx_file}")
    except Exception as e:
        print(f"Error during LibreOffice processing: {e}")
    finally:
        # Clean up temp directory
        if os.path.exists(temp_dir):
            shutil.rmtree(temp_dir)


if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description='Create a PowerPoint from a folder of images.')
    parser.add_argument('--folder',
                        type=str,
                        required=True,
                        help='Folder containing the images')
    parser.add_argument('--output',
                        type=str,
                        default=f'presentation.pptx',
                        help='Output PPTX file name')
    parser.add_argument('--no-thumbnail',
                        action='store_true',
                        help='Skip thumbnail generation (faster)')

    args = parser.parse_args()

    create_presentation_from_images(
        args.folder, 
        args.output,
        add_thumbnail=not args.no_thumbnail
    )
