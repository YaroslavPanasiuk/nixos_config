import os
from pptx import Presentation
from pdf2image import convert_from_path
from pptx.util import Inches
from images_to_pptx import process_with_libreoffice

def pdf_to_pptx(pdf_path, pptx_path):
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    images = convert_from_path(pdf_path, dpi=192)
    
    for img in images:
        temp_img_path = "temp_pdf_slide.png"
        img.save(temp_img_path, "PNG")
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        left = top = Inches(0)
        slide.shapes.add_picture(temp_img_path, left, top, 
                               width=prs.slide_width, 
                               height=prs.slide_height)
        os.remove(temp_img_path)
    
    prs.save(pptx_path)
    print(f"Successfully converted {pdf_path} to {pptx_path}")
    try:
        process_with_libreoffice(pptx_path)
    except Exception as e:
        print(f"Error generating thumbnail: {e}")

if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(description="Convert PDF to PowerPoint (PPTX)")
    parser.add_argument("--input_pdf", help="Path to input PDF file")
    parser.add_argument("--output_pptx", help="Path to output PPTX file")
    args = parser.parse_args()
    
    pdf_to_pptx(args.input_pdf, args.output_pptx)